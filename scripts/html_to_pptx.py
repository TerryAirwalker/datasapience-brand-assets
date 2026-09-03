#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Data Sapience — детерминированная сборка РЕДАКТИРУЕМОЙ презентации (Режим A).

Идея (строго по SKILL.md Шаг 3 / manifest.slides.editability):
  Автор рисует слайд как обычный HTML-макет 1920×1080 (декор + текст), помечая
  КАЖДЫЙ текстовый блок классом `t` (position:absolute;left;top;width). Движок:
    1) рендерит ФОН — тот же макет, но текст .t скрыт (color:transparent) → PNG;
    2) СНИМАЕТ из HTML для каждого .t: bbox (getBoundingClientRect) + computed style
       (шрифт, кегль, вес, цвет, выравнивание, интерлиньяж, letter-spacing, transform),
       включая per-run стили вложенных <span> (акцентные цвета/вес);
    3) СОБИРАЕТ .pptx: фон full-bleed картинкой + ЖИВЫЕ надписи python-pptx точно
       по снятым координатам; кегль pt = px×0.5, масштаб 1px = 6350 EMU.
  Шрифты Unbounded/Ping LCG встраиваются (ds_fonts.embed_fonts) — начертания
  сохраняются через отдельные typeface-имена.

CLI:
    python3 html_to_pptx.py slide1.html [slide2.html ...] -o deck.pptx
Каждый HTML = один слайд, в порядке аргументов. Итог: deck.pptx со встроенными шрифтами.
"""
import os, sys, re, json, tempfile, pathlib

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import ds_fonts

W, H = 1920, 1080
EMU = 6350           # 12192000 / 1920  (canvas 12192000×6858000 EMU = 13.333×7.5in)
PX_TO_PT = 0.5       # 1920px = 1057pt  → pt = px/2

# JS: снять все .t блоки — bbox + стиль + per-run стили + РЕАЛЬНЫЕ визуальные переносы.
# Ключ к детерминизму: авто-перенос браузера снимается через Range.getClientRects()
# и превращается в жёсткие переносы (br). В pptx перенос выключается совсем → PP не
# переливает текст по своим метрикам, строки встают ровно как в HTML.
SCRAPE_JS = r"""
() => {
  const px = v => parseFloat(v) || 0;
  const nodes = [...document.querySelectorAll('.t')];
  const out = [];
  for (const el of nodes) {
    const r = el.getBoundingClientRect();
    const cs = getComputedStyle(el);
    const runs = [];
    const transform = cs.textTransform;
    const applyTr = t => transform === 'uppercase' ? t.toUpperCase()
                        : transform === 'lowercase' ? t.toLowerCase() : t;
    const mkRun = (text, s) => ({
      text: applyTr(text),
      family: s.fontFamily, weight: s.fontWeight, style: s.fontStyle,
      color: s.color, size: px(s.fontSize),
      letter: s.letterSpacing === 'normal' ? 0 : px(s.letterSpacing),
    });
    // Текстовый узел: разбиваем на слова, по Range определяем визуальные строки,
    // на каждом переходе на новую строку вставляем br и срезаем перенос-пробел.
    const pushTextNode = (node, s) => {
      const val = node.nodeValue;
      if (!val) return;
      const toks = val.match(/\S+|\s+/g) || [];
      const range = document.createRange();
      let posn = 0, lineTop = null, cur = '';
      const flush = () => { if (cur) { runs.push(mkRun(cur, s)); cur = ''; } };
      for (const tok of toks) {
        const start = posn, end = posn + tok.length; posn = end;
        if (/^\s+$/.test(tok)) { cur += tok; continue; }
        range.setStart(node, start); range.setEnd(node, end);
        const rects = range.getClientRects();
        const top = rects.length ? Math.round(rects[0].top) : lineTop;
        if (lineTop === null) lineTop = top;
        else if (top > lineTop + 3) {           // ушли на следующую строку
          cur = cur.replace(/\s+$/, ''); flush();
          runs.push({br: true});
          lineTop = top;
        }
        cur += tok;
      }
      flush();
    };
    const walk = (node, styleEl) => {
      const s = styleEl ? getComputedStyle(styleEl) : cs;
      for (const ch of node.childNodes) {
        if (ch.nodeType === 3) pushTextNode(ch, s);
        else if (ch.nodeType === 1) {
          if (ch.tagName === 'BR') { runs.push({br: true}); continue; }
          walk(ch, ch);                          // рекурсия со стилем самого узла
        }
      }
    };
    walk(el, null);
    // нормализуем пробелы: схлопываем внутри (br не трогаем), тримим края и вокруг переносов
    for (const rn of runs) if (!rn.br) rn.text = rn.text.replace(/\s+/g, ' ');
    const isBlank = rn => !rn.br && !rn.text.trim();
    while (runs.length && isBlank(runs[0])) runs.shift();
    while (runs.length && isBlank(runs[runs.length-1])) runs.pop();
    for (let i = 0; i < runs.length; i++) {
      if (runs[i].br) continue;
      const prevBr = i === 0 || runs[i-1].br;
      const nextBr = i === runs.length-1 || runs[i+1].br;
      if (prevBr) runs[i].text = runs[i].text.replace(/^ /, '');
      if (nextBr) runs[i].text = runs[i].text.replace(/ $/, '');
    }
    const lh = cs.lineHeight === 'normal' ? px(cs.fontSize)*1.2 : px(cs.lineHeight);
    out.push({
      x: r.x, y: r.y, w: r.width, h: r.height,
      align: cs.textAlign, valign: el.dataset.valign || 'top',
      size: px(cs.fontSize), lh: lh, lhr: lh / (px(cs.fontSize) || 1),
      runs,
    });
  }
  return out;
}
"""


def _inject_fonts(html):
    """Гарантируем наличие @font-face DS в <head>, чтобы Chromium считал верные метрики."""
    if 'id="ds-fonts"' in html:
        return html
    css = ds_fonts.build_css()
    if "<head" in html:
        return re.sub(r"(<head[^>]*>)", r"\1" + css, html, count=1)
    if "<html" in html:
        return re.sub(r"(<html[^>]*>)", r"\1<head>" + css + "</head>", html, count=1)
    return css + html


def _rgba(css_color):
    """'rgb(a)(...)' -> ('RRGGBB', alpha_float)."""
    m = re.match(r"rgba?\(([^)]+)\)", css_color.strip())
    if not m:
        return "000000", 1.0
    parts = [p.strip() for p in m.group(1).split(",")]
    r, g, b = (int(round(float(parts[i]))) for i in range(3))
    a = float(parts[3]) if len(parts) > 3 else 1.0
    return f"{r:02X}{g:02X}{b:02X}", a


def _load_scrape_render(page, html_path, png_path):
    """Навигация → снять .t блоки (текст ещё виден) → скрыть текст → скриншот фона."""
    page.goto(pathlib.Path(html_path).as_uri(), wait_until="networkidle")
    page.wait_for_timeout(350)
    blocks = page.evaluate(SCRAPE_JS)
    page.add_style_tag(content=".t,.t *{color:transparent !important;text-shadow:none !important}")
    page.wait_for_timeout(120)
    page.screenshot(path=str(png_path))
    return blocks


def build(html_files, out_pptx, embed=True):
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.oxml.ns import qn

    ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
             "justify": PP_ALIGN.JUSTIFY, "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT}
    ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

    tmp = tempfile.mkdtemp(prefix="dsbuild_")
    prs = Presentation()
    prs.slide_width = Emu(W * EMU)
    prs.slide_height = Emu(H * EMU)

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        for i, hf in enumerate(html_files):
            html = _inject_fonts(pathlib.Path(hf).read_text(encoding="utf-8"))
            src_html = os.path.join(tmp, f"s{i}.html")
            pathlib.Path(src_html).write_text(html, encoding="utf-8")
            bg_png = os.path.join(tmp, f"s{i}.png")

            page = browser.new_page(viewport={"width": W, "height": H}, device_scale_factor=2)
            blocks = _load_scrape_render(page, src_html, bg_png)
            page.close()

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(bg_png, 0, 0, Emu(W * EMU), Emu(H * EMU))

            for b in blocks:
                if not b["runs"]:
                    continue
                # Переносы уже сняты как жёсткие (br) → строки встают как в HTML. Автоперенос
                # НЕ выключаем: он работает страховкой — если PP-метрики шире Chromium и
                # строка не влезает в габарит, она мягко перенесётся, а не уедет за край слайда.
                # Бокс растим с запасом по ширине (чтобы совпавшие строки не рвались раньше)
                # и якорим по выравниванию (влево / симметрично / вправо).
                pad = b["w"] * 0.08 + 12                 # запас ширины в px под метрики PP
                al = b["align"]
                if al in ("right", "end"):
                    bx = b["x"] - pad
                elif al == "center":
                    bx = b["x"] - pad / 2
                else:
                    bx = b["x"]
                box = slide.shapes.add_textbox(
                    Emu(int(bx * EMU)), Emu(int(b["y"] * EMU)),
                    Emu(int((b["w"] + pad) * EMU)), Emu(int(b["h"] * EMU)))
                tf = box.text_frame
                tf.word_wrap = True          # страховка от вылета за габарит; переносы уже жёсткие
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = ANCHOR.get(b["valign"], MSO_ANCHOR.TOP)
                line_pt = Pt(round(b["lh"] * PX_TO_PT, 1)) if b["lh"] else None
                def new_para():
                    par = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
                    par.alignment = ALIGN.get(b["align"], PP_ALIGN.LEFT)
                    par.space_before = Pt(0)
                    par.space_after = Pt(0)
                    if line_pt is not None:
                        par.line_spacing = line_pt          # точный интерлиньяж = CSS line-height
                    return par
                p = new_para()
                for rn in b["runs"]:
                    if rn.get("br"):
                        p = new_para()
                        continue
                    if rn["text"] == "":
                        continue
                    run = p.add_run(); run.text = rn["text"]
                    f = run.font
                    name, bold = ds_fonts.resolve_facename(rn["family"], rn["weight"])
                    f.name = name
                    f.bold = bold
                    f.italic = (rn["style"] == "italic")
                    f.size = Pt(round(rn["size"] * PX_TO_PT, 1))
                    hexc, alpha = _rgba(rn["color"])
                    f.color.rgb = RGBColor.from_string(hexc)
                    f.language_id = MSO_LANGUAGE_ID.RUSSIAN
                    rPr = run._r.get_or_add_rPr()
                    if rn["letter"]:
                        rPr.set("spc", str(int(round(rn["letter"] * PX_TO_PT * 100))))
                    if alpha < 0.999:
                        srgb = rPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
                        a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
                        srgb.append(a)
        browser.close()

    if embed:
        raw = os.path.join(tmp, "raw.pptx")
        prs.save(raw)
        ds_fonts.embed_fonts(raw, out_pptx)
    else:
        prs.save(out_pptx)
    return out_pptx


def main():
    args = sys.argv[1:]
    out = "deck.pptx"
    if "-o" in args:
        i = args.index("-o"); out = args[i + 1]; args = args[:i] + args[i + 2:]
    embed = "--no-embed" not in args
    args = [a for a in args if a != "--no-embed"]
    htmls = [a for a in args if a.lower().endswith((".html", ".htm"))]
    if not htmls:
        sys.exit(__doc__)
    build(htmls, out, embed=embed)
    print("saved" + (" (fonts embedded)" if embed else " (no embed — рассчитан на установленные шрифты)") + ":", out)


if __name__ == "__main__":
    main()
